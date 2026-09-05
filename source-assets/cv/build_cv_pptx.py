"""Add a PROJECTS page to the real CV deck, in the deck's own styling.

    python source-assets/cv/build_cv_pptx.py

Why edit CVs.pptx instead of drawing a new one: every font size, the Inter
family, the thin rules, the photo and the exact margins already live in that
file. Rebuilding them from scratch is how a "recreation" ends up not matching.
So this clones real shapes out of page 1 and re-fills their text -- the copy
carries the formatting with it, because in PowerPoint the styling lives on the
runs inside the shape, not in a stylesheet.

Page 1 is left alone except for one fix: Python was missing from the skills.
"""

import copy
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt

HERE = Path(__file__).parent
SRC = HERE / "CVs.pptx"
OUT = HERE / "Martinez-CV.pptx"

# --------------------------------------------------------------------------- #
# content
# --------------------------------------------------------------------------- #

# (name, tech, date, [bullets])
PROJECTS = [
    ("YODA - Your Offline Data Agent", "2026",
     "Python, Ollama, FastAPI, pandas, Agentic AI", [
         "Built a local-first data agent that answers questions over CSV and Excel files "
         "with no data leaving the machine, running a local Ollama model behind FastAPI.",
         "Designed a tool-calling loop that plans pandas operations, runs them, and "
         "checks its own output before answering.",
     ]),
    ("Aegix AI - Contract Compliance Screening", "2026",
     "Python, LangChain, Azure OpenAI, Vector Search, RAG, FastAPI", [
         "Shipped a retrieval-augmented screening tool that flags non-compliant contract "
         "clauses and cites the source paragraph behind every flag.",
         "Built the vector-search retrieval layer over Azure OpenAI embeddings, served "
         "through FastAPI with a React front end.",
     ]),
    ("Refusal Calibration - LLM Fine-Tuning", "2025",
     "PyTorch, LoRA, QLoRA, Model Evaluation, Python", [
         "Fine-tuned a base model with LoRA so it refuses out-of-scope questions instead "
         "of answering them confidently.",
         "Built a held-out evaluation set measuring over-refusal as well as refusal, so "
         "the model is not rewarded for simply declining everything.",
     ]),
    ("APAW - Dam Level and Spill-Risk Nowcaster", "2025",
     "Python, River, Online Learning, Time-Series, GitHub Actions", [
         "Deployed a self-updating forecaster that retrains on new Open-Meteo readings "
         "through a scheduled GitHub Actions job, running unattended.",
         "Used online learning so the model adapts to seasonal shift with no full retrain.",
     ]),
    ("FORGE - Fake Or Real Image Examiner", "2025",
     "PyTorch, ONNX, Computer Vision, GitHub Actions, JavaScript", [
         "Trained an image classifier separating AI-generated images from photographs, "
         "exported to ONNX so it runs in the browser with no server.",
     ]),
    ("Refresh Priority on Real Search Data", "2026",
     "scikit-learn, DuckDB, pandas, Causal Inference, Python", [
         "Ranked which pages to refresh first from real search-console data, reaching "
         "0.88 Precision@50 against a 0.86 recency baseline.",
         "Ran a leakage audit and bootstrap confidence intervals to show the gain was "
         "real and not an artefact of the split.",
     ]),
]

LINKS = "github.com/Zeref538   |   johnandreimartinez.com   |   20 projects, 18 with live demos"

# Python was absent from the original skills line, so an ATS scanning for it
# would score the CV as not having it. Everything else here is unchanged.
SKILLS_FIX = (
    "Programming & Databases: Python, SQL, PostgreSQL, MongoDB, JavaScript, Java, C, C++, React, Node.js"
)

# --------------------------------------------------------------------------- #
# layout constants, read off page 1 so page 2 lines up with it
# --------------------------------------------------------------------------- #

L = Inches(0.47)          # left margin used by every text box on page 1
RIGHT = Inches(7.75)      # right edge of the rules
TOP = Inches(0.44)        # where the name sits on page 1


# Inter 10pt across a 7.3in box fits roughly 105 characters. Without this the
# cursor advances one line per bullet, a 150-character bullet silently wraps to
# two, and the next entry lands on top of it. Nothing warns you -- python-pptx
# never renders, so the collision only shows when the file is opened.
CHARS_PER_LINE = 105
LINE = Inches(0.185)


def wrapped(lines):
    """How many rendered lines a list of strings will actually occupy."""
    return sum(max(1, -(-len(t) // CHARS_PER_LINE)) for t in lines)


def block_h(lines):
    return LINE * wrapped(lines)


def clone(shape, slide, left, top, width=None):
    """Copy a shape onto another slide, keeping all of its formatting.

    python-pptx has no "copy shape" call, so this deep-copies the underlying
    XML element and appends it to the target slide's shape tree. The copy keeps
    its fonts, sizes and colours because those attributes live inside that XML.
    """
    el = copy.deepcopy(shape._element)
    slide.shapes._spTree.append(el)
    new = slide.shapes[-1]
    new.left, new.top = left, top
    if width is not None:
        new.width = width
    return new


def set_lines(shape, lines):
    """Replace a text box's lines, reusing the first paragraph as the template.

    Each new paragraph is a copy of the original, so the font, size and bullet
    level come along. Writing plain text into a fresh paragraph instead would
    reset it to the theme default -- the usual reason a generated deck stops
    matching the one it was copied from.
    """
    tf = shape.text_frame
    template = copy.deepcopy(tf.paragraphs[0]._p)

    for p in list(tf._txBody.findall(
            "{http://schemas.openxmlformats.org/drawingml/2006/main}p")):
        tf._txBody.remove(p)

    for text in lines:
        p = copy.deepcopy(template)
        tf._txBody.append(p)
        runs = p.findall(
            "{http://schemas.openxmlformats.org/drawingml/2006/main}r")
        for extra in runs[1:]:
            p.remove(extra)
        if runs:
            runs[0].find(
                "{http://schemas.openxmlformats.org/drawingml/2006/main}t"
            ).text = text
    return shape


def main():
    prs = Presentation(SRC)
    page1 = prs.slides[0]
    by_name = {s.name: s for s in page1.shapes}

    # --- fix the skills line on page 1 ---
    skills = by_name["TextBox 45"]
    old = [t for p in skills.text_frame.paragraphs
           for t in ["".join(r.text for r in p.runs)] if t.strip()]
    old[0] = SKILLS_FIX
    set_lines(skills, old)

    # --- page 2, same blank layout as page 1 ---
    page2 = prs.slides.add_slide(page1.slide_layout)

    heading_tpl = by_name["TextBox 12"]     # "EDUCATIONAL BACKGROUND"
    rule_tpl = by_name["AutoShape 3"]       # the hairline under a heading
    title_tpl = by_name["TextBox 24"]       # bold entry name, left
    date_tpl = by_name["TextBox 23"]        # date, right
    body_tpl = by_name["TextBox 25"]        # the bullet block

    y = TOP
    set_lines(clone(heading_tpl, page2, L, y), ["PROJECTS"])
    clone(rule_tpl, page2, L, y + Inches(0.26))
    set_lines(clone(body_tpl, page2, L, y + Inches(0.32), Inches(7.3)), [LINKS])
    y += Inches(0.60)

    for name, date, tech, bullets in PROJECTS:
        set_lines(clone(title_tpl, page2, L, y, Inches(5.6)), [name])
        set_lines(clone(date_tpl, page2, Inches(5.79), y, Inches(1.97)), [date])
        # "Tech:" on its own line is what makes the stack machine-readable --
        # an ATS keyword-matches that list directly.
        set_lines(clone(body_tpl, page2, L, y + Inches(0.21), Inches(7.3)),
                  ["Tech: " + tech])
        set_lines(clone(body_tpl, page2, L, y + Inches(0.42), Inches(7.3)),
                  bullets)
        y += Inches(0.42) + block_h(bullets) + Inches(0.26)

    prs.save(OUT)

    # A page 2 that overflows the sheet silently loses its last project, so
    # check the last shape actually sits on the page.
    bottom = max(s.top + (s.height or 0) for s in page2.shapes) / 914400
    print(f"{OUT.name}: {len(prs.slides.__iter__.__self__._sldIdLst)} pages, "
          f"page 2 content ends at {bottom:.2f}in of "
          f"{prs.slide_height / 914400:.2f}in")
    assert bottom < prs.slide_height / 914400, "page 2 overflows the sheet"


# --------------------------------------------------------------------------- #
# one-page version
# --------------------------------------------------------------------------- #

# What gets cut to buy room for PROJECTS on a single sheet, and why:
#   - certificate descriptions ("Covers foundations of AI/ML...") -- nobody
#     reads them and they cost five lines
#   - the Freelance VA entry -- real work, but unrelated to ML, and it is the
#     cheapest thing on the page to lose
#   - projects drop to one bullet each
# Everything kept is flowed down the page by a cursor instead of sitting at the
# hand-placed coordinates, because once you delete a block the fixed positions
# below it leave a hole.

# Rewritten short rather than sliced: on one page every bullet has to fit a
# single rendered line (~105 characters), and the two-page wording does not.
ONE_PAGE_PROJECTS = [
    ("YODA - Your Offline Data Agent", "2026",
     "Python, Ollama, FastAPI, pandas, Agentic AI",
     ["Local-first agent answering questions over CSV/Excel with no data leaving the machine."]),
    ("Aegix AI - Contract Compliance Screening", "2026",
     "Python, LangChain, Azure OpenAI, Vector Search, RAG, FastAPI",
     ["RAG tool flagging non-compliant clauses, citing the source paragraph for every flag."]),
    ("Refusal Calibration - LLM Fine-Tuning", "2025",
     "PyTorch, LoRA, QLoRA, Model Evaluation, Python",
     ["LoRA fine-tune that refuses out-of-scope questions; evaluated for over-refusal too."]),
    ("APAW - Dam Level and Spill-Risk Nowcaster", "2025",
     "Python, River, Online Learning, Time-Series, GitHub Actions",
     ["Self-updating forecaster retraining on live Open-Meteo data via GitHub Actions."]),
]

CERTS_SHORT = [
    ("Google Advanced Data Analytics Professional Certificate", "Coursera", "July 2026"),
    ("Google AI Professional Certificate", "Coursera", "June 2026"),
    ("AWS AI Practitioner", "Udemy", "June 2026"),
    ("Building with the Claude API / Intro to Model Context Protocol", "Anthropic", "2026"),
    ("OCI Certified AI Foundations Associate", "Oracle University", "July 2025"),
    ("Google Data Analytics Professional Certificate", "Coursera", "May 2025"),
]

SKILLS_ONE_PAGE = [
    "Programming & Databases: Python, SQL, PostgreSQL, MongoDB, JavaScript, Java, C, C++, React, Node.js",
    "AI/ML & Data Science: PyTorch, scikit-learn, LLM Fine-Tuning (LoRA, QLoRA), RAG, LangChain, Agentic AI, Ollama, CNNs, Computer Vision, Time-Series",
    "Frameworks & Tools: FastAPI, Docker, Git/GitHub, GitHub Actions, ONNX, n8n, pandas, Kaggle",
    "Cloud & Visualization: Azure OpenAI, Google Cloud, Oracle Cloud, AWS, Vercel, Tableau, Power BI",
]

EDUCATION_LINES = [
    "Dean's Lister 2023-2026, with placements in Python, Database Design, and C Programming.",
    "Specialized in AI/ML Engineering: agentic AI, RAG, and LLM fine-tuning, built end-to-end.",
    "Technical Team Member, English Society; Directorate Member, JPCS.",
]

KEEP = {"Group 7", "TextBox 20", "TextBox 21", "TextBox 22", "AutoShape 2"}


def one_page():
    prs = Presentation(SRC)
    slide = prs.slides[0]
    by_name = {s.name: s for s in slide.shapes}

    heading_tpl = copy.deepcopy(by_name["TextBox 12"]._element)
    rule_tpl = copy.deepcopy(by_name["AutoShape 3"]._element)
    title_tpl = copy.deepcopy(by_name["TextBox 24"]._element)
    date_tpl = copy.deepcopy(by_name["TextBox 23"]._element)
    body_tpl = copy.deepcopy(by_name["TextBox 25"]._element)

    # add the links line to the contact block before anything moves
    set_lines(by_name["TextBox 21"],
              ["Address:", "Phone:", "Email:", "LinkedIn:", "Links:"])
    set_lines(by_name["TextBox 22"], [
        "San Jose Del Monte, Bulacan, Philippines",
        "0993-446-6975",
        "martinezjandrei8425@gmail.com",
        "/john-andrei-martinez-499a0b343",
        "github.com/Zeref538   |   johnandreimartinez.com",
    ])
    by_name["TextBox 21"].height = Inches(1.13)
    by_name["TextBox 22"].height = Inches(1.13)

    for shape in list(slide.shapes):
        if shape.name not in KEEP:
            shape._element.getparent().remove(shape._element)

    def put(tpl, left, top, width=None):
        el = copy.deepcopy(tpl)
        slide.shapes._spTree.append(el)
        sh = slide.shapes[-1]
        sh.left, sh.top = left, top
        if width is not None:
            sh.width = width
        return sh

    def section(title, y):
        set_lines(put(heading_tpl, L, y), [title])
        put(rule_tpl, L, y + Inches(0.25))
        return y + Inches(0.34)

    y = Inches(2.22)

    y = section("PROJECTS", y)
    for name, date, tech, bullets in ONE_PAGE_PROJECTS:
        set_lines(put(title_tpl, L, y, Inches(5.5)), [name])
        set_lines(put(date_tpl, Inches(5.79), y, Inches(1.97)), [date])
        set_lines(put(body_tpl, L, y + Inches(0.2), Inches(7.3)),
                  ["Tech: " + tech])
        set_lines(put(body_tpl, L, y + Inches(0.39), Inches(7.3)), bullets)
        y += Inches(0.39) + block_h(bullets) + Inches(0.13)

    y += Inches(0.06)
    y = section("EXPERIENCE", y)
    set_lines(put(title_tpl, L, y, Inches(5.5)),
              ["FlyRank AI - ML Engineering Intern"])
    set_lines(put(date_tpl, Inches(5.79), y, Inches(1.97)), ["2026 - Present"])
    exp = [
        "Prepared and validated datasets, supported model testing, automated workflows in Python.",
        "Built data pipelines and internal tools that cut manual handling from daily work.",
    ]
    set_lines(put(body_tpl, L, y + Inches(0.2), Inches(7.3)), exp)
    y += Inches(0.2) + block_h(exp) + Inches(0.13)

    y = section("EDUCATIONAL BACKGROUND", y)
    set_lines(put(title_tpl, L, y, Inches(5.5)),
              ["Our Lady of Fatima University | BS Computer Science"])
    set_lines(put(date_tpl, Inches(5.79), y, Inches(1.97)), ["2023 - Present"])
    set_lines(put(body_tpl, L, y + Inches(0.2), Inches(7.3)), EDUCATION_LINES)
    y += Inches(0.2) + block_h(EDUCATION_LINES) + Inches(0.13)

    y = section("CERTIFICATIONS", y)
    set_lines(put(body_tpl, L, y, Inches(7.3)),
              [f"{n} - {i} ({d})" for n, i, d in CERTS_SHORT])
    y += block_h([f"{n} - {i} ({d})" for n, i, d in CERTS_SHORT]) + Inches(0.13)

    y = section("SKILLS", y)
    set_lines(put(body_tpl, L, y, Inches(7.3)), SKILLS_ONE_PAGE)
    y += block_h(SKILLS_ONE_PAGE)

    out = HERE / "Martinez-Resume.pptx"
    prs.save(out)
    page_h = prs.slide_height / 914400
    print(f"{out.name}: 1 page, content ends at {y / 914400:.2f}in of {page_h:.2f}in")
    assert y / 914400 < page_h - 0.15, "one-pager overflows the sheet"


if __name__ == "__main__":
    main()
    one_page()
